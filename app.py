import streamlit as st
from docx import Document
from pptx import Presentation
import openai
import os

# --- Password Prompt (always appears) ---
password = st.text_input("Enter password to access the marking agent:", type="password")
if password != "Roxieandcleothecats1!":
    st.warning("Please enter the correct password to continue.")
    st.stop()

st.title("OCR and BTEC Marking Agent (AI-Powered)")

openai.api_key = st.secrets["OPENAI_API_KEY"] if "OPENAI_API_KEY" in st.secrets else os.getenv("OPENAI_API_KEY")

# Updated units dictionary with grading scheme
units = {
    "Unit 101 – Health and Safety in ICT": {"file": "unit101_knowledge.txt", "grading": "btec"},
    "Unit 318 – Software Installation and Upgrade": {"file": "unit318_knowledge.txt", "grading": "btec"},
    "Unit 22 – Big Data Analytics (OCR)": {"file": "unit22_knowledge.txt", "grading": "ocr"},
    "Unit 302 – Develop Own Effectiveness and Professionalism": {"file": "unit302_knowledge.txt", "grading": "btec"}
}

selected_unit = st.selectbox("Select unit to mark:", list(units.keys()))
unit_info = units[selected_unit]
knowledge = None

def load_knowledge(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

knowledge = load_knowledge(unit_info["file"])
grading_scheme = unit_info["grading"]

st.header(f"Knowledge Base for {selected_unit}")
st.text_area("Knowledge", knowledge, height=300)

# --- Clean, prominent uploader area ---
st.markdown("## Upload Your Assignment Files")
st.markdown(
    """
    Drag and drop up to two files (Word and/or PowerPoint) below:
    """,
    unsafe_allow_html=True
)

uploaded_files = st.file_uploader(
    "",
    type=["docx", "pptx"],
    accept_multiple_files=True,
    key="main_drop"
)

assignment_text = ""
if uploaded_files:
    st.success("Files uploaded!")
    for uploaded_file in uploaded_files:
        if uploaded_file.name.endswith(".docx"):
            doc = Document(uploaded_file)
            assignment_text += "\n".join([para.text for para in doc.paragraphs]) + "\n"
        elif uploaded_file.name.endswith(".pptx"):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        assignment_text += shape.text + "\n"

st.markdown("### Assignment Content Preview")
st.text_area("Assignment Text", assignment_text, height=300)

if st.button("Mark with AI"):
    with st.spinner("Marking..."):
        # Grading instructions based on scheme
        if grading_scheme == "btec":
            grading_instructions = """
- Only use Pass/Not Yet Achieved for grading. Do not mention Merit or Distinction.
- For each criterion, state if it is fully met (Pass) or not yet achieved, and use a traffic light symbol: 🟢 for fully met, 🔴 for not yet achieved.
- For criteria that are not yet achieved, provide a concise explanation of what is missing and a clear, actionable suggestion for improvement (3–4 sentences).
- For criteria that are fully met, briefly mention what is covered and any minor detail that could strengthen the evidence.
- For each learning outcome, give one short positive comment and one short suggestion for improvement (each no more than one sentence).
- Indicate if Pass is currently achievable. If not, advise what is needed for Pass.
"""
        elif grading_scheme == "ocr":
            grading_instructions = """
- Use Pass, Merit, and Distinction for grading as appropriate.
- For each criterion, state if it is fully met, partly met, or not met, and use a traffic light symbol: 🟢 for fully met, 🟡 for partly met, 🔴 for not met.
- For criteria that are not met, provide a concise explanation of what is missing and a clear, actionable suggestion for improvement (3–4 sentences).
- For criteria that are partly met, briefly state what is covered, what is missing, and how the student can fully meet the criterion (3–4 sentences).
- If a criterion is mostly met and only minor detail is missing, consider it as met for grading purposes, but mention what could be improved for a stronger grade.
- If the work is borderline between two grades, state this in your summary and explain your reasoning.
- For each learning outcome, give one short positive comment and one short suggestion for improvement (each no more than one sentence).
- Indicate which grading level is currently achievable: Pass, Merit, or Distinction.
- If some criteria are missing or incomplete, briefly advise what is needed for the next grade.
"""

        common_instructions = """
- If the assignment is incomplete or a partial submission, review only what is present and give brief guidance for improvement.
- If you suspect any content is AI-generated and not referenced, flag this and remind the student that unreferenced AI-generated work is not allowed.
- Do not rewrite or summarise the assignment; only review and comment.
- Always include the summary and grading at the very end, even if you must shorten earlier feedback to fit the token limit. If you are running out of space, prioritise the summary and grading over repeating details from earlier sections.
- Keep the entire feedback under 1200 tokens.
- Use UK English throughout.
- Use a friendly, supportive, and motivating tone. Avoid formal or academic language.
- Do not include a greeting or sign-off.
- End with a brief summary (max two sentences) stating what grade is currently achievable and encourage the student to keep improving.
"""

        improved_prompt = f"""
You are an experienced OCR and BTEC assessor for {selected_unit} (Level 3). Only use the content of the knowledge file below for all grading, definitions, and feedback. Do not use the internet, external sources, or your own assumptions.
{knowledge}
Assignment: {assignment_text}
Instructions:
{grading_instructions}
{common_instructions}
"""

        try:
            response = openai.chat.completions.create(
                model="gpt-4.1-mini",
                messages=[{"role": "user", "content": improved_prompt}],
                max_tokens=1200,
                temperature=0.2,
            )
            feedback = response.choices[0].message.content
            st.header("AI-Powered Feedback")
            st.write(feedback)
        except openai.RateLimitError:
            st.error("OpenAI rate limit reached. Please wait and try again later, or check your usage at https://platform.openai.com/usage.")
    st.info("Select your unit, upload up to two assignment files, and view the knowledge base above.")